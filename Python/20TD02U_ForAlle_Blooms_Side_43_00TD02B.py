python
from pptx import Presentation

def lag_presentasjon(tittel, slides):
    presentasjon = Presentation()
    tittel_slide = presentasjon.slides.add_slide(presentasjon.slide_layouts[0])
    tittel_slide.shapes.title.text = tittel

    for slide_tittel, slide_innhold in slides.items():
        slide = presentasjon.slides.add_slide(presentasjon.slide_layouts[1])
        slide.shapes.title.text = slide_tittel
        slide.shapes.placeholders[1].text = slide_innhold

    presentasjon.save(f"{tittel}.pptx")

slides = {
    "Introduksjon": "Dette er en introduksjon.",
    "Metode": "Metoden vi brukte var...",
    "Resultater": "Resultatene viser..."
}

lag_presentasjon("Prosjekt Presentasjon", slides)