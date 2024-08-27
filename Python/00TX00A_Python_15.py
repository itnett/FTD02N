python
from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    """
    Create a presentation for the board meeting.
    """
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_1.shapes.title
    subtitle = slide_1.placeholders[1]
    title.text = "Board Meeting Presentation"
    subtitle.text = "Arta Proff AS - January 2023"
    
    # Slide 2: SWOT Analysis
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_2.shapes.title
    title.text = "SWOT Analysis"
    content = slide_2.shapes.placeholders[1]
    content.text = (
        "Strengths:\n"
        "- Strong reputation for timely delivery\n"
        "- Competent leadership\n"
        "- Solid equity\n"
        "\nWeaknesses:\n"
        "- High operational costs\n"
        "- Dependence on current market conditions\n"
        "- Limited market diversification\n"
        "\nOpportunities:\n"
        "- Potential for product development\n"
        "- Exploring new markets\n"
        "- Technological advancements\n"
        "\nThreats:\n"
        "- Increasing costs\n"
        "- Economic downturn\n"
        "- Health and productivity issues among employees\n"
    )
    
    # Slide 3: Cost Reduction Suggestions
    slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_3.shapes.title
    title.text = "Cost Reduction Suggestions"
    content = slide_3.shapes.placeholders[1]
    content.text = (
        "- Negotiate better terms with suppliers\n"
        "- Implement energy-saving measures\n"
        "- Optimize workforce efficiency\n"
        "- Review and renegotiate service contracts\n"
        "- Reduce unnecessary overhead\n"
    )
    
    # Slide 4: Investment Proposal
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide_4.shapes.title
    title.text = "Investment Proposal"
    content = slide_4.shapes.placeholders[1]
    content.text = (
        "Description: Investing in a new automated production line\n"
        "Estimated Cost: 500,000 NOK\n"
        "Expected Savings per Year: 100,000 NOK\n"
    )
    
    prs.save("/mnt/data/Board_Meeting_Presentation.pptx")
    print("Presentation saved as 'Board_Meeting_Presentation.pptx'")

# Example usage
create_presentation()